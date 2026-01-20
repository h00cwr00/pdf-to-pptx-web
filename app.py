import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import easyocr
from PIL import Image
import io
import cv2
import numpy as np

# 設定網頁標題與寬度
st.set_page_config(page_title="NotebookLM PDF 轉 PPTX 神器", page_icon="📄", layout="centered")

# --- 快取模型，避免每次變動都重新載入 (提升速度) ---
@st.cache_resource
def load_ocr_model():
    return easyocr.Reader(['ch_tra', 'en'])

# --- 核心處理函式 (V18) ---
def repair_pdf_stream(pdf_bytes):
    """ PDF 結構修復 (針對記憶體流) """
    try:
        doc = fitz.open("pdf", pdf_bytes)
        repair_buffer = io.BytesIO()
        doc.save(repair_buffer, garbage=4, deflate=True)
        return repair_buffer.getvalue()
    except Exception:
        return pdf_bytes

def enhance_image_for_ocr(img_bgr):
    gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)
    clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
    enhanced = clahe.apply(gray)
    return cv2.cvtColor(enhanced, cv2.COLOR_GRAY2BGR)

def clean_watermark_smart(img_bgr):
    h, w = img_bgr.shape[:2]
    roi_x = int(w * 0.85)
    roi_y = int(h * 0.95)
    roi = img_bgr[roi_y:h, roi_x:w]
    if roi.size == 0: return img_bgr

    hsv = cv2.cvtColor(roi, cv2.COLOR_BGR2HSV)
    lower_gray = np.array([0, 0, 50])
    upper_gray = np.array([180, 50, 200])
    mask = cv2.inRange(hsv, lower_gray, upper_gray)
    kernel = np.ones((3, 3), np.uint8)
    mask_dilated = cv2.dilate(mask, kernel, iterations=1)
    roi_cleaned = cv2.inpaint(roi, mask_dilated, 5, cv2.INPAINT_NS)
    img_bgr_result = img_bgr.copy()
    img_bgr_result[roi_y:h, roi_x:w] = roi_cleaned
    return img_bgr_result

def convert_pdf_to_pptx(uploaded_file, progress_bar, status_text):
    # 讀取檔案流
    pdf_bytes = uploaded_file.read()
    pdf_bytes = repair_pdf_stream(pdf_bytes)
    
    doc = fitz.open("pdf", pdf_bytes)
    total_pages = len(doc)
    
    prs = Presentation()
    # 解析度設定
    zoom = 2.0 
    mat = fitz.Matrix(zoom, zoom)
    
    # 載入模型
    reader = load_ocr_model()

    for i, page in enumerate(doc):
        # 更新進度
        progress = (i + 1) / total_pages
        progress_bar.progress(progress)
        status_text.text(f"正在處理第 {i+1} / {total_pages} 頁... (AI 運算中，請稍候)")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        w_pt = Pt(page.rect.width)
        h_pt = Pt(page.rect.height)
        prs.slide_width = w_pt
        prs.slide_height = h_pt

        # 影像處理
        pix = page.get_pixmap(matrix=mat)
        img_np = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
        img_bgr = cv2.cvtColor(img_np, cv2.COLOR_RGBA2BGR if pix.n==4 else cv2.COLOR_RGB2BGR)
        
        img_bgr = clean_watermark_smart(img_bgr)
        img_enhanced = enhance_image_for_ocr(img_bgr)

        # OCR
        result = reader.readtext(img_enhanced, paragraph=True, x_ths=1.0, y_ths=0.5)
        
        mask = np.zeros(img_bgr.shape[:2], dtype=np.uint8)
        texts_to_add = []
        
        for (bbox, text) in result:
            tl = bbox[0]
            x_c, y_c = tl[0], tl[1]
            if not (x_c > img_bgr.shape[1]*0.8 and y_c > img_bgr.shape[0]*0.9):
                pt1 = (int(bbox[0][0]), int(bbox[0][1]))
                pt2 = (int(bbox[2][0]), int(bbox[2][1]))
                cv2.rectangle(mask, pt1, pt2, 255, -1)
                texts_to_add.append((bbox, text))

        # Inpainting
        kernel = np.ones((5, 5), np.uint8)
        mask_dilated = cv2.dilate(mask, kernel, iterations=1)
        img_clean = cv2.inpaint(img_bgr, mask_dilated, 5, cv2.INPAINT_NS)
        
        # 貼上背景
        img_rgb = cv2.cvtColor(img_clean, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(img_rgb)
        image_stream = io.BytesIO()
        pil_img.save(image_stream, format='PNG')
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, 0, 0, width=w_pt, height=h_pt)

        # 重建文字 (V18 智慧字體邏輯)
        for (bbox, text) in texts_to_add:
            xs = [p[0] for p in bbox]
            ys = [p[1] for p in bbox]
            x_min, y_min = min(xs), min(ys)
            w_box, h_box = max(xs) - x_min, max(ys) - y_min
            
            textbox = slide.shapes.add_textbox(Pt(x_min/zoom), Pt(y_min/zoom), Pt(w_box/zoom+20), Pt(h_box/zoom+10))
            textbox.fill.background()
            textbox.line.fill.background()
            
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            if len(text) > 10: 
                p.font.size = Pt(18)
            else:
                font_size = (h_box/zoom) * 0.8
                p.font.size = Pt(min(max(14, font_size), 44))

    # 存檔到記憶體
    output_io = io.BytesIO()
    prs.save(output_io)
    output_io.seek(0)
    return output_io

# --- 網頁介面 UI ---
st.title("📄 NotebookLM PDF 轉 PPTX")
st.markdown("### V18 線上版：去水印 + 智慧OCR")
st.info("💡 說明：上傳 PDF 後，AI 會自動進行影像增強、去字、背景修復並轉為 PPTX。")

uploaded_file = st.file_uploader("拖曳 PDF 到此處", type=["pdf"])

if uploaded_file is not None:
    if st.button("🚀 開始轉換 (Start)", type="primary"):
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            pptx_io = convert_pdf_to_pptx(uploaded_file, progress_bar, status_text)
            
            status_text.success("✅ 轉換完成！")
            progress_bar.progress(100)
            
            st.download_button(
                label="📥 下載 PPTX (Download)",
                data=pptx_io,
                file_name=f"{uploaded_file.name}_fixed.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        except Exception as e:
            st.error(f"發生錯誤：{e}")